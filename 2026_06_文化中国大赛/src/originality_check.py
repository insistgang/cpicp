#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
originality_check.py · 作品文案原创性自检(服务"原创新作"提交要求)

提交清单要求"原创新作,不得使用已公开发表/在系列赛其他赛事获奖的作品"。
本脚本复用 similarity_search 的 TF-IDF 余弦,做一个**本地、零外部依赖**的相似度自检:
  - 把本作品的若干文案段(默认取 build_pptx 的 8 页讲解备注)与一组"参照文本"
    (公开范文 / 自己历史投稿 / 队内其他稿)逐段比对,
  - 报告每段最相似的参照文本与相似度,高于阈值的标红提醒人工复核。

这**不是**抄袭判定器(只看字面/字组重合,非语义),只作"提交前自查、降低无意撞稿风险"
的提示工具;最终原创性由作者负责声明。

用法:
  from originality_check import check_originality
  report = check_originality(my_passages, reference_passages, threshold=0.5)

`python originality_check.py` 自测:构造"高度雷同/完全原创"两组,验证能正确区分。

中文 UTF-8。纯标准库 + 本目录 similarity_search(无 numpy/网络)。
"""
import os
import sys

from similarity_search import tfidf_matrix, cosine

HERE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(HERE)


def check_originality(passages, references, threshold=0.5):
    """逐段比对 passages 与 references,返回报告。

    passages:   list[str] 待检文案(本作品)
    references: list[str] 参照文本(范文/历史稿/他人稿)
    threshold:  float 余弦相似度阈值,>= 视为"需人工复核"

    返回 dict:
      items: [ {idx, text, best_ref_idx, best_ref, score, flagged} ]
      max_score, flagged_count, threshold
    若 references 为空,所有段 score=0(无可比对,视作未发现雷同)。
    """
    passages = [p for p in passages if p and p.strip()]
    references = [r for r in references if r and r.strip()]
    if not passages:
        return {"items": [], "max_score": 0.0, "flagged_count": 0,
                "threshold": threshold}

    # 把待检段与参照段放进同一 TF-IDF 空间(共享 IDF 才可比),前段是 passages。
    all_texts = passages + references
    vecs = tfidf_matrix(all_texts)
    p_vecs = vecs[:len(passages)]
    r_vecs = vecs[len(passages):]

    items = []
    for i, (ptext, pv) in enumerate(zip(passages, p_vecs)):
        best_j, best_s = -1, 0.0
        for j, rv in enumerate(r_vecs):
            s = cosine(pv, rv)
            if s > best_s:
                best_s, best_j = s, j
        items.append({
            "idx": i,
            "text": ptext,
            "best_ref_idx": best_j,
            "best_ref": references[best_j] if best_j >= 0 else None,
            "score": round(best_s, 4),
            "flagged": best_s >= threshold,
        })

    max_score = max((it["score"] for it in items), default=0.0)
    flagged = sum(1 for it in items if it["flagged"])
    return {"items": items, "max_score": max_score,
            "flagged_count": flagged, "threshold": threshold}


def format_report(report):
    """把报告渲染成可读文本。"""
    lines = []
    th = report["threshold"]
    lines.append(f"原创性自检报告 (阈值={th}; ≥阈值需人工复核)")
    lines.append(f"  待检段数: {len(report['items'])}  "
                 f"最高相似度: {report['max_score']:.3f}  "
                 f"标记复核: {report['flagged_count']} 段")
    lines.append("-" * 56)
    for it in report["items"]:
        mark = "⚠️ 复核" if it["flagged"] else "  通过"
        head = it["text"][:24] + ("…" if len(it["text"]) > 24 else "")
        lines.append(f"  [{mark}] 段{it['idx']} 相似{it['score']:.3f}  «{head}»")
        if it["flagged"] and it["best_ref"]:
            rhead = it["best_ref"][:24] + ("…" if len(it["best_ref"]) > 24 else "")
            lines.append(f"            ↳ 最像参照: «{rhead}»")
    return "\n".join(lines)


def load_pptx_notes():
    """尝试读取 output/作品_初稿.pptx 的每页备注作为待检文案;读不到则返回 []。"""
    pptx_path = os.path.join(PROJ, "output", "作品_初稿.pptx")
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        notes = []
        for sl in prs.slides:
            if sl.has_notes_slide:
                t = sl.notes_slide.notes_text_frame.text.strip()
                if t:
                    notes.append(t)
        return notes
    except Exception:
        return []


def _selftest():
    ok = True

    def check(c, m):
        nonlocal ok
        print(("  ✅ " if c else "  ❌ ") + m)
        ok = ok and c

    # 1) 高度雷同应被抓出 (近乎复制的两句字组高度重合)
    mine = [
        "见物，更见人，一部书的生命史就是普通人日常生活的精神史。",  # 与 ref0 高度雷同
        "颜之推生于南北朝乱世，历经四朝，为教子保家而著家训。",        # 与 ref1 雷同
    ]
    refs = [
        "见物，更见人，一部书的生命史正是普通人日常生活的精神史。",
        "颜之推生于南北朝乱世，历经四朝颠沛，为教子保家著成家训。",
        "春江潮水连海平，海上明月共潮生。",  # 干扰项(无关)
    ]
    rep = check_originality(mine, refs, threshold=0.5)
    check(rep["flagged_count"] == 2, f"两段高度雷同均被标记(flagged={rep['flagged_count']})")
    check(rep["items"][0]["best_ref_idx"] == 0, "段0 最像 参照0")
    check(rep["items"][1]["best_ref_idx"] == 1, "段1 最像 参照1")
    check(rep["max_score"] > 0.5, f"最高相似度>阈值(={rep['max_score']:.3f})")

    # 2) 完全原创应低于阈值
    original = ["这是一段全然不同的、关于古籍数字化方法的独立论述。"]
    rep2 = check_originality(original, refs, threshold=0.5)
    check(rep2["flagged_count"] == 0, f"原创段未被误标(flagged={rep2['flagged_count']})")
    check(rep2["max_score"] < 0.5, f"原创段相似度<阈值(={rep2['max_score']:.3f})")

    # 3) 边界: 空待检 / 空参照 不崩
    check(check_originality([], refs)["items"] == [], "空待检文案返回空报告")
    rep3 = check_originality(mine, [])
    check(rep3["max_score"] == 0.0 and rep3["flagged_count"] == 0,
          "空参照库时无雷同(max=0)")
    check(check_originality(["  ", ""], refs)["items"] == [], "全空白待检被过滤")

    # 4) report 渲染不崩且含关键字段
    txt = format_report(rep)
    check("原创性自检报告" in txt and "最高相似度" in txt, "报告可渲染")

    # 5) 自身一致性: 文本与自己比应 ~1.0 (sanity)
    rep4 = check_originality([refs[0]], [refs[0]], threshold=0.5)
    check(rep4["max_score"] > 0.99, f"文本与自身相似≈1(={rep4['max_score']:.3f})")

    # 6) 若 PPTX 存在,演示读取备注(不强制,仅打印)
    notes = load_pptx_notes()
    if notes:
        print(f"     (读取到 PPTX 备注 {len(notes)} 段, 可作真实待检文案)")

    print("\n" + ("✅ originality_check 自测通过" if ok else "❌ originality_check 自测未通过"))
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    _selftest()
