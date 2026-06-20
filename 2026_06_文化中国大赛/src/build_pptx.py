#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
build_pptx.py · 用 python-pptx 直接生成《颜氏家训》书籍生命史 8 页 16:9 PPTX

按 素材清单与PPT制作指南.md 的 8 页分镜骨架生成:
  P1 封面 / P2 选题缘起 / P3 版本流变(图) / P4 传播下沉(图) /
  P5 教化主题(词云图) / P6 互文网络(图) / P7 当代启示 / P8 结语
每页:
  - 标题 + 要点(正文)+ 配图区(数据页插 output/fig_*.png)
  - 逐字稿写入幻灯片备注(speaker notes), 作为 ≤5min 旁白脚本
全程匿名(无校名/姓名/导师/赛事标识)。

依赖: python-pptx(已装)。配图来自 render_figures.py 生成的 PNG。
若 PNG 缺失则先调用 render_figures.render_all() 生成。

输出: output/作品_初稿.pptx
`python build_pptx.py` 自测: 生成 PPTX → 读回校验 页数=8 且每页备注非空。
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 路径 ----
HERE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.dirname(HERE)
OUTDIR = os.path.join(PROJ, "output")
PPTX_PATH = os.path.join(OUTDIR, "作品_初稿.pptx")

# ---- 16:9 画布 (EMU) ----
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- 配色 ----
C_DARK = RGBColor(0x5A, 0x2D, 0x0C)
C_BROWN = RGBColor(0x8B, 0x45, 0x13)
C_CAMEL = RGBColor(0xC1, 0x9A, 0x6B)
C_CREAM = RGBColor(0xF7, 0xF3, 0xEB)
C_PAPER = RGBColor(0xFB, 0xF7, 0xEE)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY = RGBColor(0x7A, 0x6A, 0x55)
C_SEAL = RGBColor(0x9C, 0x2A, 0x1A)

# ---- 中文字体名(系统装有) ----
FONT_SERIF = "Songti SC"     # 宋体, 标题/正文古风
FONT_HEI = "Heiti SC"        # 黑体, 要点

# ============================================================
# 8 页内容 (标题 / 要点 / 配图 / 备注逐字稿)
# ============================================================
SLIDES = [
    # P1 封面
    {
        "kind": "cover",
        "title": "见物，更见人",
        "subtitle": "《颜氏家训》的书籍生命史",
        "tag": "赛题2 · 典籍与日常",
        "bg": "fig_cover.png",
        "notes": (
            "一部书，如何走过一千四百年，走进寻常人家？这就是《颜氏家训》"
            "的生命史。我们关心的，不只是它写了什么，更是它如何被一代代人"
            "抄写、刊刻、引用，沉淀进日常。这次我们用数字人文方法追踪它——"
            "既见物，也见人。"),
    },
    # P2 选题缘起
    {
        "kind": "content",
        "title": "选题缘起 · 乱世写就的第一部家训",
        "bullets": [
            "颜之推生于南北朝乱世，历经四朝，颠沛流离。",
            "他写这部家训，不为名垂青史，只为教子保家。",
            "这是中国第一部体系化家训，诞生于战火之中。",
            "核心追问：典籍中的教化智慧，如何融入日常生活？",
        ],
        "img": None,
        "notes": (
            "先说缘起。颜之推生于南北朝乱世，写这部家训，只为教子保家；它被"
            "视为中国第一部体系完备的家训。我们选它，是因为它有一个张力："
            "一部写给自家子弟的书，后来却被无数毫无血缘的家庭奉为圭臬。"
            "由此引出核心追问——教化智慧如何融入日常？"),
    },
    # P3 版本流变
    {
        "kind": "content",
        "title": "版本流变 · 越印越便宜，越传越广",
        "bullets": [
            "手抄本 → 雕版刻本 → 明清坊刻普及本 → 标点排印本 → 数字化。",
            "每一次印刷技术变革，都让这部书离普通人更近一步。",
            "数字人文方法：版本流变时间轴。",
        ],
        "img": "fig_timeline.png",
        "notes": (
            "这张时间轴整理出版本流变，七个节点：成书抄本、唐代官方著录、"
            "宋代雕版、明代坊刻普及本、清代族谱引用、民国排印，直到当代"
            "数字化。它和印刷技术的进步同步：抄本稀少，只在士族间流传；到"
            "明代坊刻已亲民到能进家塾。每一次媒介变革，都把门槛降低一截。"),
    },
    # P4 传播下沉
    {
        "kind": "content",
        "title": "传播下沉 · 从士族家学到民间日用",
        "bullets": [
            "士族家学 → 蒙学读物 → 族谱家规 → 民间日用。",
            "家训从门阀手中，走进私塾、宗族，最终融入百姓生活。",
            "这正是“大传统”走向“小传统”的转化路径。",
        ],
        "img": "fig_diffusion.png",
        "notes": (
            "这一页讲人群。传播路径是四层倒金字塔：最上是士族家学，靠手抄"
            "家传；往下是蒙学读物，与《三字经》《弟子规》并列为家塾教材；"
            "再往下是族谱家规；最底层是民间日用——理念化进婚丧嫁娶、年节"
            "礼仪。这正是“大传统”走向“小传统”的路径。"),
    },
    # P5 教化主题
    {
        "kind": "content",
        "title": "教化主题 · 数字人文里的生活智慧",
        "bullets": [
            "用 TF-IDF 字-bigram 方法，提取全书主题词。",
            "学、教、书、求益、读——高频教化关键词浮现。",
            "这些词，就是这部书教给中国人的生活智慧。",
        ],
        "img": "fig_wordcloud.png",
        "notes": (
            "这一页拉近到文本。词云用 TF-IDF 方法对样本做字与字组加权，权重"
            "越高字号越大。“学”“教”“书”这类核心字与“求益”这样的字组被推"
            "到前面——可见这部书强调的始终是怎样教、怎样学、为何读书。需说明："
            "此为演示样本，词频结论随原文补全会变，严谨结论待补全本重算。"),
    },
    # P6 互文网络
    {
        "kind": "content",
        "title": "互文网络 · 一条家训脉络的传承",
        "bullets": [
            "《颜氏家训》并非孤立存在。",
            "从《温公家范》到《朱子家训》，再到《弟子规》《三字经》。",
            "用相似检索还原一条清晰的知识传承网络。",
        ],
        "img": "fig_network.png",
        "notes": (
            "一部书的生命，也在它与别的书的关系里。这张互文网络用相似检索，"
            "以《颜氏家训》为核心，连向宋代的《温公家范》《袁氏世范》"
            "《三字经》与清代的《朱子家训》《弟子规》。强度同为方法演示，"
            "精确数值待补原文比对——但形状已经提示：典籍是在彼此引用中被"
            "记住的。"),
    },
    # P7 当代启示
    {
        "kind": "content",
        "title": "当代启示 · 千年智慧的今日回响",
        "bullets": [
            "家风家教在今天被重新提起。",
            "古籍数字化，让千年典籍触手可及。",
            "典籍中的智慧，依然在回答当代人的生活问题。",
        ],
        "img": None,
        "notes": (
            "这样一部古书，对今天有什么意义？一是家风家教被重新提起，人们"
            "重新关心该把什么价值传给下一代——这正是《颜氏家训》始终在回答"
            "的问题。二是古籍数字化：时间轴、传播路径、互文网络，都依赖典籍"
            "被结构化为可计算的数据。典籍的智慧并未过时，只是换了更易接近的"
            "方式。"),
    },
    # P8 结语
    {
        "kind": "ending",
        "title": "见物，更见人。",
        "subtitle": "一部书的生命史，就是无数普通人日常生活的精神史。",
        "bg": "fig_ending.png",
        "notes": (
            "最后回到开头那句——见物，更见人。我们追踪一部书的版本、传播与"
            "互文，表面是在看物；真正浮现的，却是无数普通人如何借这部书教"
            "孩子、立家规、过日子。一部书的生命史，归根结底就是无数普通人"
            "日常生活的精神史。谢谢。"),
    },
]


# ============================================================
# 工具函数
# ============================================================
def _blank_layout(prs):
    """取空白版式(通常 index 6)。"""
    for layout in prs.slide_layouts:
        if layout.name and "Blank" in layout.name:
            return layout
    return prs.slide_layouts[6]


def _fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _set_run(run, text, size, color, font=FONT_SERIF, bold=False):
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    # 同时设置东亚字体, 保证 PowerPoint/Keynote 都用中文字
    try:
        rPr = run._r.get_or_add_rPr()
        from pptx.oxml.ns import qn
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", font)
    except Exception:
        pass


def _add_text(slide, x, y, w, h, lines, *, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, font, bold)。每条一段。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, spec in enumerate(lines):
        text, size, color, font, bold = spec
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(8)
        run = p.add_run()
        _set_run(run, text, size, color, font, bold)
    return tb


def _img_size(path):
    """读 PNG 宽高(像素), 不依赖 PIL: 解析 IHDR。"""
    with open(path, "rb") as f:
        head = f.read(33)
    # PNG signature(8) + len(4)+'IHDR'(4) + width(4)+height(4)
    if head[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    w = int.from_bytes(head[16:20], "big")
    h = int.from_bytes(head[20:24], "big")
    return w, h


def _add_image_fit(slide, path, box_x, box_y, box_w, box_h):
    """等比缩放图片填入 box, 居中。"""
    dims = _img_size(path)
    if not dims:
        return slide.shapes.add_picture(path, box_x, box_y, width=box_w)
    iw, ih = dims
    scale = min(box_w / iw, box_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    x = box_x + (box_w - w) // 2
    y = box_y + (box_h - h) // 2
    return slide.shapes.add_picture(path, x, y, width=w, height=h)


def _set_notes(slide, text):
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    tf.text = text


# ============================================================
# 建页
# ============================================================
def _build_cover(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    bg_path = os.path.join(OUTDIR, spec["bg"])
    if os.path.exists(bg_path):
        slide.shapes.add_picture(bg_path, 0, 0, width=SLIDE_W, height=SLIDE_H)
    else:
        _fill_bg(slide, C_CREAM)
    # 左侧竖向装饰条
    _add_rect(slide, Inches(0.9), Inches(2.0), Inches(0.12), Inches(3.2), C_BROWN)
    # 标语(大)
    _add_text(slide, Inches(1.3), Inches(2.1), Inches(9.5), Inches(1.5),
              [(spec["title"], 54, C_DARK, FONT_SERIF, True)],
              align=PP_ALIGN.LEFT)
    # 副标题(作品名)
    _add_text(slide, Inches(1.32), Inches(3.6), Inches(10), Inches(1.0),
              [(spec["subtitle"], 30, C_BROWN, FONT_SERIF, False)],
              align=PP_ALIGN.LEFT)
    # 赛题标签
    _add_text(slide, Inches(1.34), Inches(4.7), Inches(8), Inches(0.6),
              [(spec["tag"], 18, C_GREY, FONT_HEI, False)],
              align=PP_ALIGN.LEFT)
    _set_notes(slide, spec["notes"])
    return slide


def _build_content(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    _fill_bg(slide, C_PAPER)
    # 顶部标题条
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), C_DARK)
    _add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.06), C_CAMEL)
    _add_text(slide, Inches(0.7), Inches(0.18), Inches(12), Inches(0.85),
              [(spec["title"], 28, C_WHITE, FONT_SERIF, True)],
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    has_img = bool(spec.get("img"))
    if has_img:
        # 左要点 + 右配图
        bullets = [(f"·  {b}", 17, C_DARK, FONT_HEI, False) for b in spec["bullets"]]
        _add_text(slide, Inches(0.7), Inches(1.7), Inches(4.6), Inches(5.2),
                  bullets, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        img_path = os.path.join(OUTDIR, spec["img"])
        _add_image_fit(slide, img_path,
                       Inches(5.5), Inches(1.5),
                       Inches(7.4), Inches(5.6))
    else:
        # 纯要点, 居中大字
        bullets = [(f"·  {b}", 24, C_DARK, FONT_HEI, False) for b in spec["bullets"]]
        _add_text(slide, Inches(1.4), Inches(2.0), Inches(10.5), Inches(4.6),
                  bullets, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    _set_notes(slide, spec["notes"])
    return slide


def _build_ending(prs, layout, spec):
    slide = prs.slides.add_slide(layout)
    bg_path = os.path.join(OUTDIR, spec["bg"])
    if os.path.exists(bg_path):
        slide.shapes.add_picture(bg_path, 0, 0, width=SLIDE_W, height=SLIDE_H)
    else:
        _fill_bg(slide, C_DARK)
    _add_text(slide, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.4),
              [(spec["title"], 60, C_CREAM, FONT_SERIF, True)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(1.0),
              [(spec["subtitle"], 22, C_CAMEL, FONT_SERIF, False)],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    _set_notes(slide, spec["notes"])
    return slide


def build(pptx_path=PPTX_PATH, ensure_figs=True):
    # 确保配图存在
    needed = [s.get("img") or s.get("bg") for s in SLIDES if s.get("img") or s.get("bg")]
    missing = [n for n in needed if not os.path.exists(os.path.join(OUTDIR, n))]
    if missing and ensure_figs:
        print(f"  配图缺失 {missing}, 调用 render_figures 生成...")
        import render_figures
        render_figures.render_all()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    layout = _blank_layout(prs)

    builders = {"cover": _build_cover, "content": _build_content, "ending": _build_ending}
    for spec in SLIDES:
        builders[spec["kind"]](prs, layout, spec)

    os.makedirs(OUTDIR, exist_ok=True)
    prs.save(pptx_path)
    print(f"  [OK] {pptx_path}")
    return pptx_path


# ============================================================
# 自测: 生成 → 读回校验 页数=8 + 备注非空
# ============================================================
def _selftest():
    ok = True

    def check(c, m):
        nonlocal ok
        print(("  ✅ " if c else "  ❌ ") + m)
        ok = ok and c

    print("=== 生成 PPTX ===")
    path = build()
    exists = os.path.exists(path)
    size = os.path.getsize(path) if exists else 0
    check(exists and size > 10000, f"PPTX 文件生成且非空 ({size} bytes)")

    print("=== 读回校验 ===")
    prs = Presentation(path)
    slides = list(prs.slides)
    check(len(slides) == 8, f"页数 = 8 (实际 {len(slides)})")

    # 16:9 画布
    check(abs(prs.slide_width - SLIDE_W) < 100 and abs(prs.slide_height - SLIDE_H) < 100,
          f"画布 16:9 ({prs.slide_width}x{prs.slide_height} EMU)")

    # 每页备注非空
    notes_ok = 0
    for i, sl in enumerate(slides, 1):
        has_notes = sl.has_notes_slide
        txt = sl.notes_slide.notes_text_frame.text.strip() if has_notes else ""
        nonempty = bool(txt)
        if nonempty:
            notes_ok += 1
        else:
            print(f"     ⚠️ 第{i}页备注为空")
    check(notes_ok == 8, f"每页备注非空 ({notes_ok}/8)")

    # 数据页含配图
    pic_pages = 0
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sl in slides:
        if any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in sl.shapes):
            pic_pages += 1
    check(pic_pages >= 4, f"含配图的页数 ≥4 (实际 {pic_pages})")

    # 备注总字数(估算讲解时长: 中文约 ~3.3 字/秒, ≤5min=300s → ≤~990字)
    total_chars = sum(
        len(sl.notes_slide.notes_text_frame.text)
        for sl in slides if sl.has_notes_slide)
    est_sec = total_chars / 3.3
    check(est_sec <= 300, f"讲解稿约 {total_chars} 字, 估时 {est_sec:.0f}s ≤300s")

    # 匿名检查: 备注/标题不含敏感词
    SENSITIVE = ["大学", "学院", "导师", "指导教师", "姓名", "学号"]
    all_text = []
    for sl in slides:
        for sh in sl.shapes:
            if sh.has_text_frame:
                all_text.append(sh.text_frame.text)
        if sl.has_notes_slide:
            all_text.append(sl.notes_slide.notes_text_frame.text)
    blob = "\n".join(all_text)
    hit = [w for w in SENSITIVE if w in blob]
    check(not hit, f"匿名检查: 无敏感词 ({'命中:' + ','.join(hit) if hit else '通过'})")

    print("\n" + ("✅ build_pptx 自测通过" if ok else "❌ build_pptx 自测未通过"))
    sys.exit(0 if ok else 1)


if __name__ == "__main__":
    _selftest()
